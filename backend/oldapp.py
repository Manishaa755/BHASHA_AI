# app.py — Full Bhasha application (single app instance)
# Run: uvicorn app:app --reload --host 127.0.0.1 --port 8000

import os
import re
import io
import base64
import asyncio
import time
import json
import mimetypes
from typing import Optional, Tuple, List, Dict, Any
from pathlib import Path

# Optional file parsers (each guarded)
try:
    import pdfplumber
except Exception:
    pdfplumber = None
try:
    import docx  # python-docx
except Exception:
    docx = None
try:
    import openpyxl
except Exception:
    openpyxl = None
try:
    from pptx import Presentation
except Exception:
    Presentation = None

# audio utils
try:
    from pydub import AudioSegment
except Exception:
    AudioSegment = None

# language detect
from langdetect import detect, DetectorFactory

# http client
import httpx

# fastapi
from fastapi import (
    FastAPI, File, UploadFile, HTTPException, Query, Body, Form, Request, Depends, Header
)
from fastapi.responses import (
    JSONResponse, FileResponse, HTMLResponse, RedirectResponse, Response, PlainTextResponse
)
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware

# env
from dotenv import load_dotenv, find_dotenv

# load env
load_dotenv(find_dotenv(), override=True)

# deterministic langdetect
DetectorFactory.seed = 42

APP_TITLE = "Bhasha — Multilingual Voice AI (STT + TTS + S2S + Lang Detect)"
APP_VERSION = "1.0.0"


def _coerce_model(value, valid, default, name):
    if not value or value.lower() not in valid:
        print(f"⚠️ {name}='{value}' invalid/unset; using '{default}'")
        return default
    return value.lower()


VALID_STT = {"whisper-1"}
VALID_CHAT = {"gpt-4o-mini", "gpt-4o"}
VALID_TTS = {"tts-1", "gpt-4o-mini-tts"}

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
OPENAI_API_BASE = os.getenv("OPENAI_API_BASE", "https://api.openai.com/v1")
WHISPER_MODEL = _coerce_model(
    os.getenv("WHISPER_MODEL"), VALID_STT, "whisper-1", "WHISPER_MODEL")
CHAT_MODEL = _coerce_model(os.getenv("CHAT_MODEL"),
                           VALID_CHAT, "gpt-4o-mini", "CHAT_MODEL")
TTS_MODEL = _coerce_model(os.getenv("TTS_MODEL"),
                          VALID_TTS, "tts-1", "TTS_MODEL")
TTS_DEFAULT_VOICE = os.getenv(
    "TTS_VOICE", os.getenv("ELEVENLABS_VOICE_ID", "alloy"))
TTS_DEFAULT_FORMAT = os.getenv("TTS_FORMAT", "mp3")

# search config
SEARCH_PROVIDER = (os.getenv("SEARCH_PROVIDER") or "").lower().strip()
BING_API_KEY = os.getenv("BING_API_KEY", "")
SERPAPI_KEY = os.getenv("SERPAPI_KEY", "")

PROJECT_ROOT = Path(__file__).resolve().parent

# Try to use env var first, then auto-detect the real frontend path
default_frontend = Path(r"D:\Superintell ai\my\frontend")
if default_frontend.is_dir():
    STATIC_DIR = str(default_frontend)
else:
    # fallback: backend/frontend
    STATIC_DIR = str(PROJECT_ROOT / "frontend")

# Allow overriding via env
STATIC_DIR = os.getenv("FRONTEND_DIR", STATIC_DIR)

# create app
app = FastAPI(title=APP_TITLE, version=APP_VERSION)

# --- optional api_keys HTTP router (only include after `app` exists) ---
try:
    import api_keys_http
    if hasattr(api_keys_http, "router"):
        app.include_router(api_keys_http.router)
        print("✅ included api_keys_http.router")
    else:
        print("⚠️ api_keys_http imported but no 'router' found")
except Exception as e:
    # safe fallback: don't crash app if the module fails to import
    print("⚠️ Could not import/include api_keys_http:", repr(e))

# Include auth and billing routers safely (they should expose `router`)
auth = None
try:
    import auth as auth_module  # type: ignore
    auth = auth_module
    if hasattr(auth_module, "router"):
        app.include_router(auth_module.router)
        print("✅ included auth.router")
    else:
        print("⚠️ auth module found but no `router` exported")
except Exception as e:
    print("⚠️ Could not import/include auth:", repr(e))
    auth = None

billing = None
try:
    import billing as billing_module  # type: ignore
    billing = billing_module
    if hasattr(billing_module, "router"):
        app.include_router(billing_module.router)
        print("✅ included billing.router")
    else:
        print("⚠️ billing module found but no `router` exported")
except Exception as e:
    print("⚠️ Could not import/include billing:", repr(e))
    billing = None

# Try to import validate_api_key_header from auth or api_keys, fallback to internal error-raising dependency.
validate_api_key_header = None
try:
    if auth and hasattr(auth, "validate_api_key_header"):
        validate_api_key_header = auth.validate_api_key_header
        print("✅ using auth.validate_api_key_header dependency")
except Exception:
    pass

if not validate_api_key_header:
    try:
        import api_keys  # backend/api_keys.py
        if hasattr(api_keys, "validate_api_key_header"):
            validate_api_key_header = api_keys.validate_api_key_header
            print("✅ using api_keys.validate_api_key_header dependency")
    except Exception as e:
        print("⚠️ Could not import api_keys.validate_api_key_header:", repr(e))

if not validate_api_key_header:
    async def validate_api_key_header(x_api_key: Optional[str] = Header(None)):
        """
        Fallback dependency: rejects requests because no key-validation provider is available.
        Replace with real implementation in auth.py or backend/api_keys.py.
        """
        raise HTTPException(
            status_code=500,
            detail="Server misconfiguration: validate_api_key_header not available. Create auth.py or backend/api_keys.py exporting validate_api_key_header."
        )

# --------------------------------------------------------------------
# CORS
# --------------------------------------------------------------------
EXPOSE = "X-Reply-Text-B64, X-Source-Lang, X-Target-Lang"
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # tighten in production
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=[h.strip() for h in EXPOSE.split(",")],
)

# --------------------------------------------------------------------
# Header-safety helpers
# --------------------------------------------------------------------
_CTL_RE = re.compile(r"[\x00-\x1F\x7F]")


def safe_header_value(value: str, max_len: int = 1024) -> str:
    if value is None:
        return ""
    value = _CTL_RE.sub("", str(value))
    value = value.encode("ascii", "ignore").decode("ascii")
    return value[:max_len]


def b64_header(value: str, max_len: int = 8192) -> str:
    if not value:
        return ""
    b = base64.b64encode(value.encode("utf-8")).decode("ascii")
    return safe_header_value(b, max_len=max_len)

# --------------------------------------------------------------------
# Last-mile CORS + safety
# --------------------------------------------------------------------


@app.middleware("http")
async def ensure_cors_header(request: Request, call_next):
    try:
        resp = await call_next(request)
    except Exception as e:
        print("❌ Unhandled exception:", repr(e))
        resp = JSONResponse(status_code=500, content={
                            "error": "internal_server_error", "detail": str(e)[:800]})
    if "access-control-allow-origin" not in (k.lower() for k in resp.headers.keys()):
        resp.headers["Access-Control-Allow-Origin"] = "*"
        resp.headers["Access-Control-Allow-Headers"] = "*"
        resp.headers["Access-Control-Allow-Methods"] = "*"
        resp.headers["Access-Control-Expose-Headers"] = EXPOSE
    return resp

# --------------------------------------------------------------------
# Optional fastText and CLD3
# --------------------------------------------------------------------
FASTTEXT_MODEL_PATH = os.getenv("FASTTEXT_MODEL_PATH", "")
_ft = None
if FASTTEXT_MODEL_PATH:
    try:
        import fasttext
        if os.path.isfile(FASTTEXT_MODEL_PATH):
            _ft = fasttext.load_model(FASTTEXT_MODEL_PATH)
        else:
            print(f"⚠️ FASTTEXT_MODEL_PATH not found: {FASTTEXT_MODEL_PATH}")
    except Exception as e:
        print(f"⚠️ fastText unavailable ({e}).")

_cld3 = None
try:
    import pycld3  # optional
    _cld3 = pycld3.NNetLanguageIdentifier(min_num_bytes=0, max_num_bytes=2048)
except Exception as e:
    print(f"⚠️ pycld3 unavailable ({e}).")
try:
    import demo_proxy
    if hasattr(demo_proxy, "router"):
        app.include_router(demo_proxy.router)
        print("✅ included demo_proxy.router")
except Exception as e:
    print("⚠️ Could not include demo_proxy:", e)

# --------------------------------------------------------------------
# Web search utilities
# --------------------------------------------------------------------


def needs_web(q: str) -> bool:
    ql = (q or "").lower().strip()
    triggers = (
        "current", "today", "latest", "price", "prices", "rate", "interest", "weather", "score", "scores",
        "breaking", "stock", "exchange", "convert", "conversion", "when is", "deadline", "holiday",
        "who won", "live", "update", "news", "market", "value", "forecast", "schedule", "fixture",
        "result", "results", "rd rate", "fd rate", "interest rate", "inflation", "gdp", "population",
        "time now", "now", "right now"
    )
    if any(w in ql for w in triggers):
        return True
    if ("?" in ql) or re.search(r"\b(202\d|20[3-9]\d|today|tomorrow|yesterday)\b", ql):
        return True
    return False


def looks_offline_reply(text: str) -> bool:
    if not text:
        return False
    t = text.lower()
    patterns = (
        "i'm unable to provide real-time data",
        "i cannot browse",
        "i can’t browse",
        "i cannot access the internet",
        "real-time information",
        "up-to-date information",
        "check the official website",
        "visit the official website",
        "as an ai model i don't have real-time",
        "as an ai language model i don't have",
    )
    return any(p in t for p in patterns)


async def web_search(query: str, num: int = 5) -> List[Dict]:
    results: List[Dict] = []
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            if SEARCH_PROVIDER == "bing" and BING_API_KEY:
                url = "https://api.bing.microsoft.com/v7.0/search"
                params = {"q": query, "count": num,
                          "mkt": "en-IN", "safeSearch": "Moderate"}
                headers = {"Ocp-Apim-Subscription-Key": BING_API_KEY}
                r = await client.get(url, params=params, headers=headers)
                if r.is_success:
                    data = r.json()
                    for item in (data.get("webPages", {}).get("value", []) or [])[:num]:
                        results.append({"title": item.get("name", ""), "url": item.get(
                            "url", ""), "snippet": item.get("snippet", "")})
            elif SEARCH_PROVIDER == "serpapi" and SERPAPI_KEY:
                url = "https://serpapi.com/search.json"
                params = {"engine": "google", "q": query,
                          "num": num, "hl": "en", "api_key": SERPAPI_KEY}
                r = await client.get(url, params=params)
                if r.is_success:
                    data = r.json()
                    for item in (data.get("organic_results") or [])[:num]:
                        results.append({"title": item.get("title", ""), "url": item.get(
                            "link", ""), "snippet": item.get("snippet", "")})
    except Exception as e:
        print("⚠️ web_search error:", repr(e))
    return results

_TAG_RE = re.compile(r"<(script|style)[\s\S]*?</\1>|<[^>]+>", re.IGNORECASE)
_WS_RE = re.compile(r"\s+")


async def fetch_and_extract(url: str, max_chars: int = 4000) -> str:
    try:
        async with httpx.AsyncClient(timeout=20, follow_redirects=True) as client:
            r = await client.get(url, headers={"User-Agent": "Mozilla/5.0"})
            if not r.is_success:
                return ""
            text = r.text or ""
            text = _TAG_RE.sub(" ", text)
            text = _WS_RE.sub(" ", text).strip()
            return text[:max_chars]
    except Exception as e:
        print("⚠️ fetch_and_extract error for", url, ":", repr(e))
        return ""


async def answer_with_web(query: str, target_lang: str = "en") -> Dict:
    hits = await web_search(query, num=5)
    if not hits:
        return {"answer": "", "sources": []}
    extracts = []
    for h in hits[:3]:
        body = await fetch_and_extract(h["url"])
        if not body:
            continue
        extracts.append(
            f"### Source: {h['title']}\nURL: {h['url']}\n\n{body}\n")
    context = "\n\n".join(extracts)[:12000]

    system = (
        "You answer questions using ONLY the provided web extracts. "
        "Be concise (3–6 sentences), include key numbers/dates if present, and answer in the requested language code. "
        "End with a short 'Sources:' list of the top 1–3 URLs you used. If the info is missing, say so."
    )
    messages = [
        {"role": "system", "content": system},
        {"role": "user", "content": f"Language code to answer in: {target_lang}"},
        {"role": "user", "content": f"Question: {query}\n\nWeb Extracts:\n{context}"}
    ]
    answer = await openai_chat(messages, temperature=0.1)
    return {"answer": answer, "sources": [h["url"] for h in hits[:3]]}

# --------------------------------------------------------------------
# Language utilities
# --------------------------------------------------------------------


def _normalize_label(code: str) -> str:
    if not code:
        return ""
    c = code.strip().lower()
    if "-" in c or "_" in c:
        c = re.split(r"[-_]", c)[0]
    if len(c) > 3:
        c = c[:3]
    return c


def detect_language_fasttext_topk(text: str, k: int = 5) -> List[Dict]:
    out: List[Dict] = []
    if not _ft:
        return out
    t = (text or "").strip()
    if not t:
        return out
    try:
        labels, probs = _ft.predict(t.replace("\n", " "), k=k)
        for lab, p in zip(labels, probs):
            raw = lab.replace("__label__", "")
            out.append(
                {"code_raw": raw, "code_base": _normalize_label(raw), "prob": float(p)})
    except Exception:
        pass
    return out


def detect_language_cld3_topk(text: str, k: int = 5) -> List[Dict]:
    out: List[Dict] = []
    if not _cld3:
        return out
    t = (text or "").strip()
    if not t:
        return out
    try:
        langs = _cld3.FindTopNMostFreqLangs(t, k)
        for item in (langs or []):
            code = (item.language or "").lower()
            prob = float(item.probability or 0.0)
            out.append(
                {"code_raw": code, "code_base": _normalize_label(code), "prob": prob})
    except Exception:
        pass
    return out


def detect_language_langdetect(text: str) -> str:
    try:
        return detect(text) or ""
    except Exception:
        return ""

# --------------------------------------------------------------------
# OpenAI wrappers (chat, transcribe, tts) using httpx
# --------------------------------------------------------------------


async def openai_chat(user_messages: List[Dict], temperature=0.2) -> str:
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=500, detail="Missing OPENAI_API_KEY.")
    url = f"{OPENAI_API_BASE.rstrip('/')}/chat/completions"
    headers = {"Authorization": f"Bearer {OPENAI_API_KEY}",
               "Content-Type": "application/json"}
    payload = {"model": CHAT_MODEL,
               "temperature": temperature, "messages": user_messages}
    max_tries, backoff = 4, 1.0
    async with httpx.AsyncClient(timeout=180) as client:
        for attempt in range(1, max_tries + 1):
            resp = await client.post(url, headers=headers, json=payload)
            if resp.status_code == 429 or 500 <= resp.status_code < 600:
                if attempt == max_tries:
                    break
                delay = float(resp.headers.get("retry-after") or backoff)
                await asyncio.sleep(delay)
                backoff *= 2
                continue
            if not resp.is_success:
                raise HTTPException(
                    status_code=resp.status_code, detail=resp.text)
            data = resp.json()
            return (data.get("choices", [{}])[0].get("message", {}).get("content", "")).strip()
    raise HTTPException(
        status_code=502, detail="Chat completion failed after retries.")


async def openai_chat_reply(user_text: str, lang_hint: Optional[str], enforce_lang: bool = True) -> str:
    base = ("You are a multilingual assistant. Be concise (2–4 sentences). "
            "If unsure, ask a clarifying question; do not guess.")
    if enforce_lang and lang_hint and lang_hint != "auto":
        system_msg = f"{base} Always reply exclusively in language code '{lang_hint}'."
    else:
        system_msg = f"{base} Choose the most appropriate language for the user."
    return await openai_chat([{"role": "system", "content": system_msg}, {"role": "user", "content": user_text}], temperature=0.2)


async def openai_transcribe(file_bytes: bytes, filename: str, language: Optional[str] = None) -> dict:
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=500, detail="Missing OPENAI_API_KEY.")
    url = f"{OPENAI_API_BASE.rstrip('/')}/audio/transcriptions"
    headers = {"Authorization": f"Bearer {OPENAI_API_KEY}"}

    def make_files():
        files = {"file": (filename or "audio.webm", file_bytes,
                          "application/octet-stream"), "model": (None, WHISPER_MODEL)}
        if language:
            files["language"] = (None, language)
        return files

    max_tries, backoff = 4, 1.0
    async with httpx.AsyncClient(timeout=180) as client:
        for attempt in range(1, max_tries + 1):
            resp = await client.post(url, headers=headers, files=make_files())
            if resp.status_code == 400 and attempt == 1:
                # Try converting to WAV 16k mono
                try:
                    wav = to_wav_16k_mono(file_bytes)
                    files = {
                        "file": ("converted.wav", wav, "audio/wav"), "model": (None, WHISPER_MODEL)}
                    if language:
                        files["language"] = (None, language)
                    resp = await client.post(url, headers=headers, files=files)
                except Exception:
                    pass
            if resp.status_code == 429 or 500 <= resp.status_code < 600:
                if attempt == max_tries:
                    break
                delay = float(resp.headers.get("retry-after") or backoff)
                await asyncio.sleep(delay)
                backoff *= 2
                continue
            if not resp.is_success:
                raise HTTPException(
                    status_code=resp.status_code, detail=resp.text)
            data = resp.json()
            return {"text": (data.get("text") or "").strip(), "language": data.get("language")}
    raise HTTPException(
        status_code=502, detail="Transcription failed after retries.")


async def openai_dictionary_translate_to_english(term: str) -> str:
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=500, detail="Missing OPENAI_API_KEY.")
    payload = {"model": CHAT_MODEL, "temperature": 0.0, "messages": [
               {"role": "system", "content": "You are a bilingual dictionary. Return ONLY the most common English equivalent. No punctuation."},
               {"role": "user", "content": term}]}
    url = f"{OPENAI_API_BASE.rstrip('/')}/chat/completions"
    headers = {"Authorization": f"Bearer {OPENAI_API_KEY}",
               "Content-Type": "application/json"}
    async with httpx.AsyncClient(timeout=60) as client:
        resp = await client.post(url, headers=headers, json=payload)
        if not resp.is_success:
            raise HTTPException(status_code=resp.status_code, detail=resp.text)
        data = resp.json()
        translation = (data.get("choices", [{}])[0].get(
            "message", {}).get("content", "")).strip()
        return translation.split("\n")[0].strip(" '\".,;:()")


def _audio_content_type(fmt: str) -> str:
    fmt = (fmt or "").lower()
    return {
        "mp3": "audio/mpeg",
        "wav": "audio/wav",
        "opus": "audio/ogg",
        "aac": "audio/aac",
        "flac": "audio/flac",
    }.get(fmt, "application/octet-stream")


def to_wav_16k_mono(raw_bytes: bytes) -> bytes:
    if not AudioSegment:
        raise RuntimeError("pydub not available for audio conversion")
    from io import BytesIO
    buf_in = BytesIO(raw_bytes)
    audio = AudioSegment.from_file(buf_in)
    audio = audio.set_frame_rate(16000).set_channels(1)
    buf_out = BytesIO()
    audio.export(buf_out, format="wav")
    return buf_out.getvalue()

# --------------------------------------------------------------------
# LLM-powered language helpers
# --------------------------------------------------------------------


async def llm_detect_language(text: str) -> Dict:
    if not OPENAI_API_KEY:
        return {}
    system = (
        "Identify the language of the given text across 200+ languages. "
        'Return ONLY JSON: {"code":"<iso>","name":"<language name>","confidence":<0..1>} '
        "Use ISO 639-1 if available else ISO 639-3 for the code."
    )
    try:
        out = await openai_chat([{"role": "system", "content": system}, {"role": "user", "content": text}], temperature=0.0)
        try:
            obj = json.loads(out)
        except Exception:
            obj = {}
            m = re.search(r'"code"\s*:\s*"([^"]+)"', out)
            n = re.search(r'"name"\s*:\s*"([^"]+)"', out)
            c = re.search(r'"confidence"\s*:\s*([0-9.]+)', out)
            if m:
                obj["code"] = m.group(1)
            if n:
                obj["name"] = n.group(1)
            if c:
                obj["confidence"] = float(c.group(1))
        code = (obj.get("code") or "").strip().lower()
        name = (obj.get("name") or "").strip()
        conf = float(obj.get("confidence") or 0.85)
        if not code:
            return {}
        return {"code_raw": code, "code_base": _normalize_label(code), "name": name, "prob": conf}
    except Exception:
        return {}


async def llm_infer_target_language(user_text: str, default_code: str) -> str:
    if not OPENAI_API_KEY:
        return default_code
    system = (
        "Decide the requested reply language for the user text. "
        f"If a language is specified, return only its ISO code. Otherwise return only '{default_code}'."
    )
    try:
        code = await openai_chat([{"role": "system", "content": system}, {"role": "user", "content": user_text}], temperature=0.0)
        return _normalize_label(code)
    except Exception:
        return default_code


async def llm_detect_dictionary_intent(text: str) -> Dict:
    if not OPENAI_API_KEY:
        return {"intent": "none"}
    system = (
        "Detect dictionary intent like 'X ko English me kya kehte hain?'. "
        'Return ONLY JSON. If yes: {"intent":"dict","term":"..."} , else {"intent":"none"}.'
    )
    try:
        content = await openai_chat([{"role": "system", "content": system}, {"role": "user", "content": text}], temperature=0.0)
        try:
            obj = json.loads(content)
        except Exception:
            obj = {"intent": "none"}
        if obj.get("intent") == "dict" and obj.get("term"):
            return {"intent": "dict", "term": str(obj["term"]).strip()}
        return {"intent": "none"}
    except Exception:
        return {"intent": "none"}


async def ensemble_best_code(text: str) -> Dict:
    ft = detect_language_fasttext_topk(text, k=5)
    cd = detect_language_cld3_topk(text, k=5)
    ld_raw = detect_language_langdetect(text)
    ld_base = _normalize_label(ld_raw) if ld_raw else ""
    llm = await llm_detect_language(text)

    scores: Dict[str, float] = {}
    for e in ft:
        scores[e["code_base"]] = scores.get(e["code_base"], 0.0) + e["prob"]
    for e in cd:
        scores[e["code_base"]] = scores.get(e["code_base"], 0.0) + e["prob"]
    if ld_base:
        scores[ld_base] = scores.get(ld_base, 0.0) + 0.60
    if llm:
        scores[llm["code_base"]] = scores.get(
            llm["code_base"], 0.0) + float(llm.get("prob", 0.85))

    best_code, best_score = "", -1.0
    for k, v in scores.items():
        if v > best_score:
            best_code, best_score = k, v

    return {
        "best": {"code_base": best_code, "score": best_score},
        "scores": scores,
        "detectors": {
            "fasttext":  {"available": bool(_ft),  "top": ft},
            "cld3":      {"available": bool(_cld3), "top": cd},
            "langdetect": {"available": True, "code_raw": ld_raw, "code_base": ld_base},
            "llm":       {"available": bool(OPENAI_API_KEY), "top": llm},
        }
    }

# --------------------------------------------------------------------
# Health
# --------------------------------------------------------------------


@app.get("/health")
async def health():
    return {
        "status": "ok",
        "stt_model": WHISPER_MODEL,
        "chat_model": CHAT_MODEL,
        "tts_model": TTS_MODEL,
        "detectors": {
            "fasttext": bool(_ft),
            "cld3": bool(_cld3),
            "langdetect": True,
            "llm": bool(OPENAI_API_KEY),
        },
        "frontend_dir": STATIC_DIR,
        "has_api_key": bool(OPENAI_API_KEY),
        "search": {
            "provider": SEARCH_PROVIDER,
            "has_bing_key": bool(BING_API_KEY),
            "has_serpapi_key": bool(SERPAPI_KEY),
        }
    }

# --------------------------------------------------------------------
# Document helpers (extract/summarize/qa)
# --------------------------------------------------------------------
GENERIC_SUMMARY_REGEX = re.compile(
    r"^\s*(summarize( this)?|make (a )?summary|give (me )?(a )?summary|tl;dr)\s*\.?\s*$",
    flags=re.IGNORECASE
)


def _guess_ext(filename: str) -> str:
    if not filename:
        return ""
    return os.path.splitext(filename)[1].lower()


async def _extract_text_from_upload(upload: UploadFile) -> str:
    raw = await upload.read()
    name = upload.filename or ""
    ext = _guess_ext(name)
    ctype = upload.content_type or mimetypes.guess_type(name)[0] or ""

    def _as_text_bytes(b: bytes) -> str:
        try:
            return b.decode("utf-8")
        except Exception:
            return b.decode("utf-8", errors="ignore")

    # PDF
    if ext == ".pdf" or "pdf" in ctype:
        if not pdfplumber:
            return _as_text_bytes(raw)
        text = []
        with pdfplumber.open(io.BytesIO(raw)) as pdf:
            for page in pdf.pages:
                t = page.extract_text() or ""
                if t.strip():
                    text.append(t)
        return "\n\n".join(text).strip()

    # DOCX
    if ext == ".docx" or "officedocument.wordprocessingml" in ctype:
        if not docx:
            return _as_text_bytes(raw)
        d = docx.Document(io.BytesIO(raw))
        return "\n".join(p.text for p in d.paragraphs).strip()

    # PPTX
    if ext == ".pptx" or "officedocument.presentationml" in ctype:
        if not Presentation:
            return _as_text_bytes(raw)
        prs = Presentation(io.BytesIO(raw))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = (shape.text or "").strip()
                    if t:
                        lines.append(t)
        return "\n".join(lines).strip()

    # XLSX
    if ext == ".xlsx" or "officedocument.spreadsheetml" in ctype:
        if not openpyxl:
            return _as_text_bytes(raw)
        wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True)
        lines = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                vals = [str(v) for v in row if v is not None]
                if vals:
                    lines.append("\t".join(vals))
        return "\n".join(lines).strip()

    # Plain text, csv, md, json, etc.
    return _as_text_bytes(raw).strip()


async def _summarize_text(body: str, target_lang: str) -> str:
    prompt = ("Summarize the following content clearly and concisely. "
              "Return bullet points followed by a 2–3 sentence abstract.\n\n"
              f"{body[:120000]}")
    return await openai_chat_reply(prompt, target_lang or "en", enforce_lang=True)


async def _qa_over_text(body: str, question: str, target_lang: str) -> str:
    prompt = ("Answer the user’s question using ONLY the information in the provided document. "
              "If the answer isn’t present, say you can’t find it.\n\n"
              "### Document:\n"
              f"{body[:120000]}\n\n"
              f"### Question:\n{question}")
    return await openai_chat_reply(prompt, target_lang or "en", enforce_lang=True)

# --------------------------------------------------------------------
# Endpoints: /v1/reason-doc
# --------------------------------------------------------------------


@app.post("/v1/reason-doc")
async def reason_doc(
    task: str = Form("auto"),
    question: Optional[str] = Form(None),
    target_lang: str = Form("auto"),
    text: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None)
):
    try:
        doc_text = ""
        if file and file.filename:
            doc_text = await _extract_text_from_upload(file)
            if not doc_text.strip():
                raise HTTPException(
                    status_code=400, detail="Could not extract readable text from the uploaded file.")
        elif text and text.strip():
            doc_text = text.strip()
        else:
            raise HTTPException(
                status_code=400, detail="Provide either a file or some text to analyze.")

        inferred_task = task.strip().lower()
        if inferred_task == "auto":
            if question and not GENERIC_SUMMARY_REGEX.match(question or ""):
                inferred_task = "qa"
            else:
                inferred_task = "summarize"
        if question and GENERIC_SUMMARY_REGEX.match(question):
            inferred_task = "summarize"

        if target_lang == "auto":
            ens = await ensemble_best_code(doc_text[:2000])
            target_lang = ens["best"]["code_base"] or "en"
        else:
            target_lang = _normalize_label(target_lang)

        if inferred_task == "qa":
            if not question or GENERIC_SUMMARY_REGEX.match(question):
                summary = await _summarize_text(doc_text, target_lang)
                return {"task": "summarize", "target_lang": target_lang, "result": summary}
            answer = await _qa_over_text(doc_text, question, target_lang)
            return {"task": "qa", "target_lang": target_lang, "question": question, "result": answer}

        summary = await _summarize_text(doc_text, target_lang)
        return {"task": "summarize", "target_lang": target_lang, "result": summary}

    except HTTPException as he:
        return JSONResponse(status_code=he.status_code, content={"detail": he.detail})
    except Exception as e:
        print("❌ /v1/reason-doc error:", repr(e))
        return JSONResponse(status_code=500, content={"error": "internal_server_error", "detail": str(e)[:800]})

# --------------------------------------------------------------------
# Language detection mapping and endpoint
# --------------------------------------------------------------------
LANG_NAME_TO_CODE = {
    "english": "en", "hindi": "hi", "gujarati": "gu", "bengali": "bn", "bangla": "bn", "tamil": "ta", "telugu": "te",
    "marathi": "mr", "kannada": "kn", "malayalam": "ml", "punjabi": "pa", "urdu": "ur", "odia": "or", "assamese": "as",
    "nepali": "ne", "sinhala": "si", "sindhi": "sd", "french": "fr", "spanish": "es", "german": "de", "italian": "it",
    "portuguese": "pt", "russian": "ru", "arabic": "ar", "turkish": "tr", "chinese": "zh", "japanese": "ja", "korean": "ko",
    "vietnamese": "vi", "thai": "th", "indonesian": "id", "malay": "ms", "swahili": "sw", "persian": "fa", "hebrew": "he",
    # Endonyms
    "हिंदी": "hi", "ગુજરાતી": "gu", "বাংলা": "bn", "தமிழ்": "ta", "తెలుగు": "te", "मराठी": "mr", "ಕನ್ನಡ": "kn", "മലയാളം": "ml",
    "ਪੰਜਾਬੀ": "pa", "اردو": "ur", "தமிழில்": "ta", "हिन्दी": "hi", "ગુજરાતીમાં": "gu"
}


def guess_requested_lang(user_text: str) -> Optional[str]:
    t = (user_text or "").lower()
    m = re.search(
        r"\b(?:in|into)\s+([A-Za-z\u00C0-\u024F\u0400-\u04FF\u0590-\u05FF\u0600-\u06FF\u0900-\u0D7F]+)\b", t)
    if not m:
        m = re.search(
            r"\b([A-Za-z\u00C0-\u024F\u0400-\u04FF\u0590-\u05FF\u0600-\u06FF\u0900-\u0D7F]+)\s+language\b", t)
    if m:
        name = m.group(1).strip()
        code = LANG_NAME_TO_CODE.get(name)
        if code:
            return code
    return None

# --------------------------------------------------------------------
# STT chat (audio -> text reply)
# --------------------------------------------------------------------


@app.post("/chat/")
async def chat_with_bot(
    file: UploadFile = File(...),
    lang: str = Query(
        "auto", description="Language hint (e.g., hi, bn, ta-IN) or 'auto'"),
    live: bool = Query(False, description="Force web search for live info")
):
    try:
        if not file or not file.filename:
            raise HTTPException(
                status_code=400, detail="No audio file provided.")
        audio_bytes = await file.read()
        if not audio_bytes:
            raise HTTPException(
                status_code=400, detail="Uploaded file is empty.")

        lang_hint = None if lang.lower() == "auto" else lang
        stt = await openai_transcribe(audio_bytes, file.filename, lang_hint)
        user_text = stt.get("text", "")

        ensemble = await ensemble_best_code(user_text)
        src_lang = ensemble["best"]["code_base"] or (
            stt.get("language") or "auto")

        forced = guess_requested_lang(user_text)
        if forced:
            tgt_lang = forced
        else:
            tgt_lang = await llm_infer_target_language(user_text, src_lang)

        # Dictionary short path
        d = await llm_detect_dictionary_intent(user_text)
        if d.get("intent") == "dict":
            english = await openai_dictionary_translate_to_english(d["term"])
            reply = await openai_chat_reply(
                f"Say in {tgt_lang} that the '{d['term']}' is called '{english}' in English, succinctly.",
                tgt_lang, enforce_lang=True
            )
        else:
            # Live web if requested or needed
            if live or needs_web(user_text):
                web = await answer_with_web(user_text, target_lang=tgt_lang if tgt_lang != "auto" else "en")
                if web.get("answer"):
                    reply = web["answer"]
                else:
                    reply = await openai_chat_reply(user_text, tgt_lang, enforce_lang=True)
            else:
                reply = await openai_chat_reply(user_text, tgt_lang, enforce_lang=True)

        return JSONResponse({
            "user_text": user_text,
            "bot_reply": reply,
            "source_language": src_lang,
            "target_language": tgt_lang,
            "lang_detect": ensemble
        })
    except HTTPException as he:
        return JSONResponse(status_code=he.status_code, content={"detail": he.detail})
    except Exception as e:
        print("❌ /chat error:", repr(e))
        return JSONResponse(status_code=500, content={"error": "internal_server_error", "detail": str(e)[:800]})

# --------------------------------------------------------------------
# TTS (text -> speech)
# --------------------------------------------------------------------


@app.options("/v1/tts-speak")
async def tts_speak_options():
    return PlainTextResponse("", status_code=204, headers={
        "Access-Control-Allow-Origin": "*",
        "Access-Control-Allow-Headers": "*",
        "Access-Control-Allow-Methods": "*",
        "Access-Control-Expose-Headers": EXPOSE,
    })


@app.post("/v1/tts-speak")
async def tts_speak(payload: dict = Body(...), live: bool = Query(False, description="Force web search for live info")):
    try:
        txt = payload.get("text")
        user_text = (txt.strip() if isinstance(txt, str) else "")
        if not user_text:
            raise HTTPException(status_code=400, detail="Missing 'text'.")

        ensemble = await ensemble_best_code(user_text)
        src = ensemble["best"]["code_base"] or "auto"

        forced = guess_requested_lang(user_text)
        if forced:
            tgt_lang = forced
        else:
            tgt_lang = await llm_infer_target_language(user_text, src)

        voice = (payload.get("voice") or TTS_DEFAULT_VOICE).strip()
        fmt = (payload.get("format") or TTS_DEFAULT_FORMAT).lower().strip()
        tts_model = (payload.get("tts_model") or TTS_MODEL).strip()

        d = await llm_detect_dictionary_intent(user_text)
        if d.get("intent") == "dict":
            english = await openai_dictionary_translate_to_english(d["term"])
            reply_text = await openai_chat_reply(
                f"Say in {tgt_lang} that the '{d['term']}' is called '{english}' in English, succinctly.",
                tgt_lang, enforce_lang=True
            )
        else:
            if live or needs_web(user_text):
                web = await answer_with_web(user_text, target_lang=tgt_lang if tgt_lang != "auto" else "en")
                if web.get("answer"):
                    reply_text = web["answer"]
                else:
                    reply_text = await openai_chat_reply(user_text, tgt_lang, enforce_lang=True)
            else:
                reply_text = await openai_chat_reply(user_text, tgt_lang, enforce_lang=True)

        audio_bytes = await openai_tts(reply_text=reply_text, voice=voice, fmt=fmt, model=tts_model)
        media_type = _audio_content_type(fmt)

        headers = {
            "X-Reply-Text-B64": b64_header(reply_text),
            "X-Source-Lang": safe_header_value(src, 32),
            "X-Target-Lang": safe_header_value(tgt_lang, 32),
            "Cache-Control": "no-store",
            "Content-Disposition": safe_header_value(f'inline; filename="reply.{fmt}"'),
        }
        return Response(content=audio_bytes, media_type=media_type, headers=headers)
    except HTTPException as he:
        return JSONResponse(status_code=he.status_code, content={"detail": he.detail})
    except Exception as e:
        print("❌ /v1/tts-speak error:", repr(e))
        return JSONResponse(status_code=500, content={"error": "internal_server_error", "detail": str(e)[:800]})

# --------------------------------------------------------------------
# Language detection endpoint
# --------------------------------------------------------------------


@app.post("/v1/lang-detect")
async def lang_detect(text: Optional[str] = Form(None), file: Optional[UploadFile] = File(None)):
    try:
        detected_text = None
        if file and file.filename:
            audio_bytes = await file.read()
            if not audio_bytes:
                raise HTTPException(
                    status_code=400, detail="Uploaded audio is empty.")
            stt = await openai_transcribe(audio_bytes, file.filename, None)
            detected_text = (stt.get("text") or "").strip()
        elif text:
            detected_text = text.strip()

        if not detected_text:
            raise HTTPException(status_code=400, detail="No input provided.")

        ensemble = await ensemble_best_code(detected_text)

        combined: Dict[str, float] = {}
        for e in ensemble["detectors"]["fasttext"]["top"]:
            combined[e["code_base"]] = combined.get(
                e["code_base"], 0.0) + e["prob"]
        for e in ensemble["detectors"]["cld3"]["top"]:
            combined[e["code_base"]] = combined.get(
                e["code_base"], 0.0) + e["prob"]
        ld_base = ensemble["detectors"]["langdetect"]["code_base"]
        if ld_base:
            combined[ld_base] = combined.get(ld_base, 0.0) + 0.60
        llm = ensemble["detectors"]["llm"]
        if llm and llm.get("top", {}).get("code_base"):
            combined[llm["top"]["code_base"]] = combined.get(
                llm["top"]["code_base"], 0.0) + float(llm["top"].get("prob", 0.85))

        top_combined = sorted([{"code_base": k, "score": v} for k, v in combined.items(
        )], key=lambda x: x["score"], reverse=True)[:5]

        return JSONResponse({"input": detected_text, "best": ensemble["best"], "top_combined": top_combined, "detectors": ensemble["detectors"]})
    except HTTPException as he:
        return JSONResponse(status_code=he.status_code, content={"detail": he.detail})
    except Exception as e:
        print("❌ /v1/lang-detect error:", repr(e))
        return JSONResponse(status_code=500, content={"error": "internal_server_error", "detail": str(e)[:800]})

# --------------------------------------------------------------------
# Speech-to-Speech (S2S)
# --------------------------------------------------------------------


@app.post("/v1/s2s")
async def speech_to_speech(
    file: UploadFile = File(...,
                            description="Audio file from MediaRecorder or upload"),
    target_lang: Optional[str] = Query(
        None, description="(Optional) 'same'/'auto'/language code override"),
    voice: Optional[str] = Query(None, description="TTS voice"),
    format: Optional[str] = Query(
        None, description="Audio format: mp3|wav|opus|aac|flac"),
    tts_model: Optional[str] = Query(None, description="TTS model"),
    live: bool = Query(
        False, description="Force web search for live info (e.g., ?live=1)")
):
    try:
        if not file or not file.filename:
            raise HTTPException(
                status_code=400, detail="No audio file provided.")
        audio_bytes = await file.read()
        if not audio_bytes:
            raise HTTPException(
                status_code=400, detail="Uploaded file is empty.")

        # 1) STT
        stt = await openai_transcribe(audio_bytes, file.filename, None)
        user_text = (stt.get("text") or "").strip()
        if not user_text:
            raise HTTPException(
                status_code=400, detail="Could not understand audio.")

        # 2) Detect src lang
        ens = await ensemble_best_code(user_text)
        src_lang = ens["best"]["code_base"] or (stt.get("language") or "auto")

        # 3) target lang selection
        forced = guess_requested_lang(user_text)
        if forced:
            tgt_lang, enforce = forced, True
        elif target_lang:
            tl = target_lang.lower().strip()
            if tl == "same":
                tgt_lang, enforce = src_lang, True
            elif tl == "auto":
                tgt_lang = await llm_infer_target_language(user_text, src_lang)
                enforce = (tgt_lang != "auto")
            else:
                tgt_lang, enforce = _normalize_label(tl), True
        else:
            tgt_lang = await llm_infer_target_language(user_text, src_lang)
            enforce = (tgt_lang != "auto")

        # 4) Compose reply (dictionary intent short path)
        d = await llm_detect_dictionary_intent(user_text)
        if d.get("intent") == "dict":
            english = await openai_dictionary_translate_to_english(d["term"])
            reply_text = await openai_chat_reply(
                f"Say in {tgt_lang if tgt_lang != 'auto' else 'the best fitting language'} that the '{d['term']}' is called '{english}' in English, succinctly.",
                tgt_lang, enforce_lang=enforce
            )
        else:
            tried_web = False
            reply_text = ""

            if live or needs_web(user_text):
                tried_web = True
                web = await answer_with_web(user_text, target_lang=tgt_lang if tgt_lang != "auto" else "en")
                reply_text = (web.get("answer") or "").strip()

            if not reply_text:
                reply_text = await openai_chat_reply(user_text, tgt_lang, enforce_lang=enforce)
                if not tried_web and looks_offline_reply(reply_text):
                    web = await answer_with_web(user_text, target_lang=tgt_lang if tgt_lang != "auto" else "en")
                    if web.get("answer"):
                        reply_text = web["answer"]

        # 5) TTS
        v = (voice or TTS_DEFAULT_VOICE).strip()
        fmt = (format or TTS_DEFAULT_FORMAT).lower().strip()
        ttsm = (tts_model or TTS_MODEL).strip()
        audio_bytes_out = await openai_tts(reply_text=reply_text, voice=v, fmt=fmt, model=ttsm)
        media_type = _audio_content_type(fmt)

        headers = {
            "X-Reply-Text-B64": b64_header(reply_text),
            "X-Source-Lang": safe_header_value(src_lang, 32),
            "X-Target-Lang": safe_header_value(tgt_lang, 32),
            "Cache-Control": "no-store",
            "Content-Disposition": safe_header_value(f'inline; filename="reply.{fmt}"'),
        }
        return Response(content=audio_bytes_out, media_type=media_type, headers=headers)

    except HTTPException as he:
        return JSONResponse(status_code=he.status_code, content={"detail": he.detail})
    except Exception as e:
        print("❌ /v1/s2s error:", repr(e))
        return JSONResponse(status_code=500, content={"error": "internal_server_error", "detail": str(e)[:800]})

# --------------------------------------------------------------------
# -------------------- NEW: /tool/* wrappers protected by X-API-Key --------------------
# --------------------------------------------------------------------
# All endpoints below require the API key header dependency:
#   headers: { "X-API-Key": "<key>" }
# The dependency (validate_api_key_header) validates key, applies per-minute quota,
# and returns api key metadata which we accept as "api_key_info"


@app.post("/tool/stt")
async def tool_stt(file: UploadFile = File(...), api_key_info: dict = Depends(validate_api_key_header)):
    """
    STT for third-party clients. Expects X-API-Key header (validated).
    Returns JSON: { text: "...", language: "en" }
    """
    try:
        if billing:
            try:
                # Charge in INR (1.0 INR for STT)
                billing.charge_for_usage(api_key_info.get("user_id") or api_key_info.get(
                    "id") or api_key_info.get("email"), 1.0, "tool_stt")
            except HTTPException as bhe:
                return JSONResponse(status_code=bhe.status_code, content={"detail": bhe.detail})

        if not file or not file.filename:
            raise HTTPException(
                status_code=400, detail="No audio file provided.")
        audio_bytes = await file.read()
        if not audio_bytes:
            raise HTTPException(
                status_code=400, detail="Uploaded file is empty.")
        # we pass None for language to let model auto-detect
        stt = await openai_transcribe(audio_bytes, file.filename, None)
        return JSONResponse({"text": stt.get("text", ""), "language": stt.get("language")})
    except HTTPException as he:
        return JSONResponse(status_code=he.status_code, content={"detail": he.detail})
    except Exception as e:
        print("❌ /tool/stt error:", repr(e))
        return JSONResponse(status_code=500, content={"error": "internal_server_error", "detail": str(e)[:800]})


@app.post("/tool/tts")
async def tool_tts(payload: dict = Body(...), api_key_info: dict = Depends(validate_api_key_header)):
    """
    TTS wrapper for third-party clients protected by X-API-Key.
    Body (JSON): { text: "...", voice: "alloy", format: "mp3", tts_model: "tts-1" }
    Returns audio bytes with appropriate headers.
    """
    try:
        if billing:
            try:
                # Charge in INR (2.0 INR for TTS)
                billing.charge_for_usage(api_key_info.get("user_id") or api_key_info.get(
                    "id") or api_key_info.get("email"), 2.0, "tool_tts")
            except HTTPException as bhe:
                return JSONResponse(status_code=bhe.status_code, content={"detail": bhe.detail})

        txt = payload.get("text")
        user_text = (txt.strip() if isinstance(txt, str) else "")
        if not user_text:
            raise HTTPException(status_code=400, detail="Missing 'text'.")
        # choose voice/format/model
        voice = (payload.get("voice") or TTS_DEFAULT_VOICE).strip()
        fmt = (payload.get("format") or TTS_DEFAULT_FORMAT).lower().strip()
        tts_model = (payload.get("tts_model") or TTS_MODEL).strip()
        # infer languages similarly to existing tts endpoint
        ens = await ensemble_best_code(user_text)
        src = ens["best"]["code_base"] or "auto"
        forced = guess_requested_lang(user_text)
        if forced:
            tgt_lang = forced
        else:
            tgt_lang = await llm_infer_target_language(user_text, src)
        # Compose reply (dictionary short-path)
        d = await llm_detect_dictionary_intent(user_text)
        if d.get("intent") == "dict":
            english = await openai_dictionary_translate_to_english(d["term"])
            reply_text = await openai_chat_reply(
                f"Say in {tgt_lang} that the '{d['term']}' is called '{english}' in English, succinctly.",
                tgt_lang, enforce_lang=True
            )
        else:
            reply_text = await openai_chat_reply(user_text, tgt_lang, enforce_lang=True)
        audio_bytes = await openai_tts(reply_text=reply_text, voice=voice, fmt=fmt, model=tts_model)
        media_type = _audio_content_type(fmt)
        headers = {
            "X-Reply-Text-B64": b64_header(reply_text),
            "X-Source-Lang": safe_header_value(src, 32),
            "X-Target-Lang": safe_header_value(tgt_lang, 32),
            "Cache-Control": "no-store",
            "Content-Disposition": safe_header_value(f'inline; filename="reply.{fmt}"'),
        }
        return Response(content=audio_bytes, media_type=media_type, headers=headers)
    except HTTPException as he:
        return JSONResponse(status_code=he.status_code, content={"detail": he.detail})
    except Exception as e:
        print("❌ /tool/tts error:", repr(e))
        return JSONResponse(status_code=500, content={"error": "internal_server_error", "detail": str(e)[:800]})


@app.post("/tool/s2s")
async def tool_s2s(
    file: UploadFile = File(...),
    target_lang: Optional[str] = Query(None),
    voice: Optional[str] = Query(None),
    format: Optional[str] = Query(None),
    tts_model: Optional[str] = Query(None),
    live: bool = Query(False),
    api_key_info: dict = Depends(validate_api_key_header)
):
    """
    Speech-to-speech wrapper that accepts an audio file and returns TTS audio.
    Protected by X-API-Key.
    """
    try:
        if billing:
            try:
                # Charge in INR (3.0 INR for S2S)
                billing.charge_for_usage(api_key_info.get("user_id") or api_key_info.get(
                    "id") or api_key_info.get("email"), 3.0, "tool_s2s")
            except HTTPException as bhe:
                return JSONResponse(status_code=bhe.status_code, content={"detail": bhe.detail})

        if not file or not file.filename:
            raise HTTPException(
                status_code=400, detail="No audio file provided.")
        audio_bytes = await file.read()
        if not audio_bytes:
            raise HTTPException(
                status_code=400, detail="Uploaded file is empty.")
        # 1) STT
        stt = await openai_transcribe(audio_bytes, file.filename, None)
        user_text = (stt.get("text") or "").strip()
        if not user_text:
            raise HTTPException(
                status_code=400, detail="Could not understand audio.")
        # 2) detect src lang
        ens = await ensemble_best_code(user_text)
        src_lang = ens["best"]["code_base"] or (stt.get("language") or "auto")
        # 3) target selection
        forced = guess_requested_lang(user_text)
        if forced:
            tgt_lang, enforce = forced, True
        elif target_lang:
            tl = target_lang.lower().strip()
            if tl == "same":
                tgt_lang, enforce = src_lang, True
            elif tl == "auto":
                tgt_lang = await llm_infer_target_language(user_text, src_lang)
                enforce = (tgt_lang != "auto")
            else:
                tgt_lang, enforce = _normalize_label(tl), True
        else:
            tgt_lang = await llm_infer_target_language(user_text, src_lang)
            enforce = (tgt_lang != "auto")
        # 4) reply
        d = await llm_detect_dictionary_intent(user_text)
        if d.get("intent") == "dict":
            english = await openai_dictionary_translate_to_english(d["term"])
            reply_text = await openai_chat_reply(
                f"Say in {tgt_lang if tgt_lang != 'auto' else 'the best fitting language'} that the '{d['term']}' is called '{english}' in English, succinctly.",
                tgt_lang, enforce_lang=enforce
            )
        else:
            tried_web = False
            reply_text = ""
            if live or needs_web(user_text):
                tried_web = True
                web = await answer_with_web(user_text, target_lang=tgt_lang if tgt_lang != "auto" else "en")
                reply_text = (web.get("answer") or "").strip()
            if not reply_text:
                reply_text = await openai_chat_reply(user_text, tgt_lang, enforce_lang=enforce)
                if not tried_web and looks_offline_reply(reply_text):
                    web = await answer_with_web(user_text, target_lang=tgt_lang if tgt_lang != "auto" else "en")
                    if web.get("answer"):
                        reply_text = web["answer"]
        # 5) TTS
        v = (voice or TTS_DEFAULT_VOICE).strip()
        fmt = (format or TTS_DEFAULT_FORMAT).lower().strip()
        ttsm = (tts_model or TTS_MODEL).strip()
        audio_bytes_out = await openai_tts(reply_text=reply_text, voice=v, fmt=fmt, model=ttsm)
        media_type = _audio_content_type(fmt)
        headers = {
            "X-Reply-Text-B64": b64_header(reply_text),
            "X-Source-Lang": safe_header_value(src_lang, 32),
            "X-Target-Lang": safe_header_value(tgt_lang, 32),
            "Cache-Control": "no-store",
            "Content-Disposition": safe_header_value(f'inline; filename="reply.{fmt}"'),
        }
        return Response(content=audio_bytes_out, media_type=media_type, headers=headers)
    except HTTPException as he:
        return JSONResponse(status_code=he.status_code, content={"detail": he.detail})
    except Exception as e:
        print("❌ /tool/s2s error:", repr(e))
        return JSONResponse(status_code=500, content={"error": "internal_server_error", "detail": str(e)[:800]})

# tool/reason-doc and tool/lang-detect reuse existing handlers


@app.post("/tool/reason-doc")
async def tool_reason_doc(
    task: str = Form("auto"),
    question: Optional[str] = Form(None),
    target_lang: str = Form("auto"),
    text: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
    api_key_info: dict = Depends(validate_api_key_header)
):
    try:
        result = await reason_doc(task=task, question=question, target_lang=target_lang, text=text, file=file)
        if isinstance(result, JSONResponse):
            return result
        return JSONResponse(result)
    except HTTPException as he:
        return JSONResponse(status_code=he.status_code, content={"detail": he.detail})
    except Exception as e:
        print("❌ /tool/reason-doc error:", repr(e))
        return JSONResponse(status_code=500, content={"error": "internal_server_error", "detail": str(e)[:800]})


@app.post("/tool/lang-detect")
async def tool_lang_detect(text: Optional[str] = Form(None), file: Optional[UploadFile] = File(None), api_key_info: dict = Depends(validate_api_key_header)):
    try:
        result = await lang_detect(text=text, file=file)
        if isinstance(result, JSONResponse):
            return result
        return JSONResponse(result)
    except HTTPException as he:
        return JSONResponse(status_code=he.status_code, content={"detail": he.detail})
    except Exception as e:
        print("❌ /tool/lang-detect error:", repr(e))
        return JSONResponse(status_code=500, content={"error": "internal_server_error", "detail": str(e)[:800]})

# --------------------------------------------------------------------
# Static frontend (landing page + redirects)
# --------------------------------------------------------------------
if not os.path.isdir(STATIC_DIR):
    print(f"⚠️ FRONTEND_DIR not found: {STATIC_DIR}")

app.mount("/static", StaticFiles(directory=STATIC_DIR), name="static")
app.mount("/frontend/static", StaticFiles(directory=STATIC_DIR),
          name="frontend_static")


@app.get("/", response_class=FileResponse)
async def home():
    index_path = os.path.join(STATIC_DIR, "index.html")
    if not os.path.isfile(index_path):
        return HTMLResponse(f"<h1>Missing index.html</h1><p>Expected at: <code>{index_path}</code></p>", status_code=200)
    return FileResponse(index_path)


@app.get("/reasoning")
async def reasoning_redirect():
    return RedirectResponse(url="/static/reasoning.html", status_code=307)


@app.get("/stt")
async def stt_redirect():
    return RedirectResponse(url="/static/stt.html", status_code=307)


@app.get("/tts")
async def tts_redirect():
    return RedirectResponse(url="/static/tts.html", status_code=307)


@app.get("/s2s")
async def s2s_redirect():
    return RedirectResponse(url="/static/sts.html", status_code=307)


@app.get("/langdetect")
async def langdetect_redirect():
    return RedirectResponse(url="/static/langdetect.html", status_code=307)


@app.get("/sts")
async def sts_redirect():
    return RedirectResponse(url="/static/sts.html", status_code=307)


@app.get("/profile")
async def profile_redirect():
    return RedirectResponse(url="/static/profile.html", status_code=307)


@app.get("/dashboard")
async def dashboard_redirect():
    return RedirectResponse(url="/static/dashboard.html", status_code=307)


@app.get("/guide", include_in_schema=False)
async def guide_redirect():
    path = os.path.join(STATIC_DIR, "docs.html")
    if os.path.isfile(path):
        return RedirectResponse(url="/static/docs.html", status_code=307)
    return HTMLResponse(f"<h1>docs.html not found</h1><p>Expected at: <code>{path}</code></p>", status_code=404)

# --------------------------------------------------------------------
# NEW FUNCTION: openai_tts using ElevenLabs (or fallback)
# --------------------------------------------------------------------


async def openai_tts(reply_text: str, voice: str = None, fmt: str = "mp3", model: str = None) -> bytes:
    """
    Produces speech audio bytes for the given reply_text.
    Primary implementation: ElevenLabs text-to-speech API.
    If ELEVENLABS_API_KEY is not set, returns a tiny silent WAV placeholder.
    """
    eleven_api_key = os.getenv("ELEVENLABS_API_KEY") or os.getenv(
        "ELEVENLABS_KEY") or None
    eleven_voice = voice or os.getenv(
        "ELEVENLABS_VOICE_ID") or TTS_DEFAULT_VOICE
    if not eleven_api_key:
        # No ElevenLabs configured — return a tiny silent WAV as a safe fallback.
        print("⚠️ ELEVENLABS_API_KEY not configured; returning silent placeholder audio.")
        # Minimal WAV header for 0-second silent audio (very small). Better to fail in production.
        return base64.b64decode("UklGRiQAAABXQVZFZm10IBAAAAABAAEAQB8AAIA+AAACABAAZGF0YQAAAAA=")

    url = f"https://api.elevenlabs.io/v1/text-to-speech/{eleven_voice}"
    headers = {
        "xi-api-key": eleven_api_key,
        "Content-Type": "application/json",
        "Accept": "audio/mpeg, audio/wav, */*"
    }
    body = {
        "text": reply_text,
    }
    accept_map = {
        "mp3": "audio/mpeg",
        "wav": "audio/wav",
        "ogg": "audio/ogg",
        "aac": "audio/aac",
        "flac": "audio/flac"
    }
    accept_hdr = accept_map.get((fmt or "").lower(), "audio/mpeg")
    headers["Accept"] = accept_hdr

    async with httpx.AsyncClient(timeout=120) as client:
        try:
            r = await client.post(url, headers=headers, json=body)
            if r.status_code != 200:
                try:
                    err = r.json()
                except Exception:
                    err = r.text
                raise HTTPException(
                    status_code=502, detail=f"ElevenLabs TTS error: {err}")
            return r.content
        except httpx.RequestError as exc:
            print("⚠️ ElevenLabs request error:", exc)
            raise HTTPException(
                status_code=502, detail="TTS provider request failed")
        except HTTPException:
            raise
        except Exception as e:
            print("⚠️ ElevenLabs unexpected error:", repr(e))
            raise HTTPException(
                status_code=502, detail="TTS processing failed")

# --------------------------------------------------------------------
# End of file
# --------------------------------------------------------------------
